import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox
import os
import sys
import jieba
import math
import warnings
import pdfplumber
# import threading
import pickle
import json
import hashlib
from collections import defaultdict, Counter
from concurrent.futures import ThreadPoolExecutor
from PIL import Image

try:
    from docx import Document

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pandas as pd

    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

try:
    from bs4 import BeautifulSoup

    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False

def get_resource_path(relative_path):
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)

# 基础配置
warnings.filterwarnings("ignore", message=".*pkg_resources.*")
jieba.setLogLevel(jieba.logging.ERROR)

ctk.set_appearance_mode("Light")
ctk.set_default_color_theme("blue")

# 核心配色
COLORS = {
    "bg_main": "#FAEDD1", "bg_sidebar": "#1387C0",
    "text_sidebar": "#FFFFFF", "text_main": "#1387C0",
    "text_content": "#2B3A42", "text_highlight": "#F4520D",
    "btn_primary": "#F4520D", "btn_primary_hover": "#FF6E2B",

    "btn_white_bg": "#FFFFFF", "btn_white_text": "#1387C0",
    "btn_white_hover": "#F0F0F0",

    "card_bg": "#FFFFFF", "card_hover": "#E1F5FE", "card_border": "#1387C0",
    "progress": "#F4520D", "progress_bg": "#0B5A85",

    "input_bg": "#FFFFFF", "input_text": "#1387C0", "input_border": "#1387C0",

    "dropdown_fg": "#FFFFFF", "dropdown_hover": "#E1F5FE", "dropdown_text": "#1387C0",
    "header_bg": "#FFFFFF", "header_text": "#1387C0",
}

# 搜索引擎
class RankedSearchEngine:
    def __init__(self, stop_words_file='stopwords.txt'):
        self.documents = {}
        self.doc_paths = {}
        self.doc_titles = {}
        self.doc_freq = defaultdict(int)
        self.doc_term_freqs = {}
        self.total_docs = 0
        real_path = get_resource_path(stop_words_file)
        self.stop_words = self._load_stop_words(real_path)
        self.indexed_folder = ""

    def _load_stop_words(self, file_path):
        loaded = set()
        if os.path.exists(file_path):
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    loaded = {line.strip() for line in f if line.strip()}
            except:
                pass
        return loaded

    def save_index_to_disk(self, file_path):
        try:
            with open(file_path, 'wb') as f:
                pickle.dump(self.__dict__, f)
        except Exception as e:
            print(f"保存索引失败: {e}")

    def load_index_from_disk(self, file_path):
        try:
            if not os.path.exists(file_path): return False
            with open(file_path, 'rb') as f:
                tmp_dict = pickle.load(f)
                self.__dict__.update(tmp_dict)
            return True
        except:
            return False

    def add_document(self, doc_id, text, file_path, title):
        self.documents[doc_id] = text
        self.doc_paths[doc_id] = file_path
        self.doc_titles[doc_id] = title
        self.total_docs += 1

        words = jieba.lcut(text)
        # 过滤
        clean_words = [w for w in words if w not in self.stop_words and len(w.strip()) > 0]
        # 记录词频信息
        self.doc_term_freqs[doc_id] = {'counts': Counter(clean_words), 'length': len(clean_words)}
        # 更新文档频率 (DF)
        for word in set(clean_words):
            self.doc_freq[word] += 1

    def _calculate_score(self, query_words, doc_id):
        score = 0.0
        doc_data = self.doc_term_freqs[doc_id]
        if doc_data['length'] == 0: return 0.0

        for word in query_words:
            term_count = doc_data['counts'].get(word, 0)
            if term_count == 0: continue

            # TF 饱和度处理
            tf = (term_count * 2.0) / (term_count + 1.5)
            # IDF 平滑处理
            idf = math.log10(self.total_docs / (self.doc_freq.get(word, 0) + 1)) + 1.0
            score += tf * idf
        return score

    def search(self, query, top_k=20):
        # 预处理
        raw_words = jieba.lcut(query)
        query_words = [w for w in raw_words if w not in self.stop_words and len(w.strip()) > 0]
        if not query_words and len(query.strip()) > 0:
            query_words = [query.strip()]
        if not query_words: return []

        # 评分
        temp_results = []
        max_raw_score = 0
        for doc_id in self.documents:
            s = self._calculate_score(query_words, doc_id)
            if s > 0:
                temp_results.append((doc_id, s))
                if s > max_raw_score: max_raw_score = s

        # 显式降序排序
        temp_results.sort(key=lambda x: x[1], reverse=True)

        results = []
        for doc_id, s in temp_results[:top_k]:
            display_score = int((s / max_raw_score) * 99) if max_raw_score > 0 else 0
            results.append({
                'score': display_score,
                'title': self.doc_titles[doc_id],
                'path': self.doc_paths[doc_id],
                'content': self.documents[doc_id]
            })
        return results

    def _extract_content(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        content = ""
        try:
            if ext == '.pdf':
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        txt = page.extract_text()
                        if txt: content += txt + "\n"
            elif ext == '.docx' and HAS_DOCX:
                doc = Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs])
            elif ext == '.pptx' and HAS_PPTX:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): content += shape.text + "\n"
            elif ext in ['.xlsx', '.xls', '.csv'] and HAS_PANDAS:
                if ext == '.csv':
                    df = pd.read_csv(file_path, on_bad_lines='skip')
                else:
                    df = pd.read_excel(file_path)
                content = df.to_string()
            elif ext in ['.html', '.xml'] and HAS_BS4:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    soup = BeautifulSoup(f, 'html.parser')
                    content = soup.get_text()
            elif ext in ['.txt', '.md', '.py', '.json', '.log', '.xml']:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    try:
                        with open(file_path, 'r', encoding='gbk') as f:
                            content = f.read()
                    except:
                        pass
        except Exception as e:
            print(f"❌ 解析失败: {os.path.basename(file_path)} -> {e}")
        return content

# 界面逻辑
class VibrantSearchApp(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("本地搜索")
        self.geometry("1100x750")
        self.configure(fg_color=COLORS["bg_main"])

        self.engine = RankedSearchEngine()
        self.current_folder = ""
        self.folder_history = []
        self.search_history = []

        self.executor = ThreadPoolExecutor(max_workers=1)
        self.indexes_dir = "indexes"
        if not os.path.exists(self.indexes_dir): os.makedirs(self.indexes_dir)

        # 布局
        self.sidebar_frame = ctk.CTkFrame(self, width=260, corner_radius=0, fg_color=COLORS["bg_sidebar"])
        self.sidebar_frame.pack(side="left", fill="y", expand=False)

        self.main_container = ctk.CTkFrame(self, corner_radius=0, fg_color="transparent")
        self.main_container.pack(side="right", fill="both", expand=True)

        self.content_layer = ctk.CTkFrame(self.main_container, fg_color="transparent")
        self._setup_background()

        # 左控件
        self.logo_label = ctk.CTkLabel(
            self.sidebar_frame,
            text="文档索引",
            font=ctk.CTkFont(family="微软雅黑", size=32, weight="bold"),
            text_color="white",
            fg_color="transparent"
        )
        self.logo_label.pack(padx=20, pady=(40, 30))

        ctk.CTkLabel(self.sidebar_frame, text="当前索引库:", font=ctk.CTkFont(size=12), text_color="#B3E5FC",
                     anchor="w").pack(padx=25, pady=(0, 5), fill="x")

        self.folder_var = ctk.StringVar(value="请选择或添加...")

        # 文件夹选择
        self.folder_menu = ctk.CTkOptionMenu(
            self.sidebar_frame,
            variable=self.folder_var,
            height=40,
            fg_color="white",
            text_color=COLORS["text_main"],
            button_color="#0B5A85",
            button_hover_color="#084463",
            dropdown_fg_color="white",
            dropdown_text_color=COLORS["text_main"],
            dropdown_hover_color="#E1F5FE",
            font=ctk.CTkFont(size=13, weight="bold"),
            command=self.on_folder_change
        )
        self.folder_menu.pack(padx=25, pady=(0, 15), fill="x")

        self.btn_add_folder = ctk.CTkButton(self.sidebar_frame, text="+ 添加新文件夹", command=self.browse_new_folder,
                                            height=35,
                                            fg_color="transparent", border_width=1, border_color="white",
                                            text_color="white", hover_color="#0B5A85")
        self.btn_add_folder.pack(padx=25, pady=(0, 20), fill="x")

        self.btn_index = ctk.CTkButton(self.sidebar_frame, text="⚡ 重建当前库索引", command=self.start_indexing,
                                       state="disabled", height=45,
                                       fg_color=COLORS["btn_primary"], hover_color=COLORS["btn_primary_hover"],
                                       text_color="white", font=ctk.CTkFont(size=14, weight="bold"))
        self.btn_index.pack(padx=25, pady=10, fill="x")

        tk.Frame(self.sidebar_frame, height=1, bg="white").pack(fill="x", padx=30, pady=30)

        self.status_label = ctk.CTkLabel(self.sidebar_frame, text="就绪", font=ctk.CTkFont(size=12),
                                         text_color="#E1F5FE", wraplength=220, anchor="w")
        self.status_label.pack(padx=25, pady=(0, 10), fill="x")

        self.progress = ctk.CTkProgressBar(self.sidebar_frame, mode="indeterminate", progress_color=COLORS["progress"],
                                           fg_color=COLORS["progress_bg"])

        # 右控件
        ctk.CTkLabel(self.content_layer, text="你想搜索些什么?",
                     font=ctk.CTkFont(family="Arial", size=36, weight="bold"),
                     text_color=COLORS["text_main"]).pack(anchor="w", pady=(0, 20), padx=40)

        search_frame = ctk.CTkFrame(self.content_layer, fg_color="transparent")
        search_frame.pack(fill="x", pady=(0, 20), padx=40)

        self.search_combo = ctk.CTkComboBox(
            search_frame,
            height=55,
            font=ctk.CTkFont(size=18),
            border_width=2,
            corner_radius=10,
            fg_color="white",
            border_color=COLORS["input_border"],
            text_color=COLORS["input_text"],
            button_color=COLORS["bg_sidebar"],
            button_hover_color="#0B5A85",
            dropdown_fg_color="white",
            dropdown_text_color=COLORS["text_main"],
            dropdown_hover_color=COLORS["dropdown_hover"]
        )
        self.search_combo.set("")
        self.search_combo.pack(side="left", fill="x", expand=True)
        self.search_combo._entry.bind('<Return>', lambda event: self.perform_search_async())

        # 垃圾桶
        self.btn_clear_history = ctk.CTkButton(
            search_frame,
            text="🗑️",
            width=50, height=55,
            fg_color="white",
            border_width=2,
            border_color="#C62828",
            hover_color="#FFEBEE",
            text_color="#C62828",
            font=ctk.CTkFont(size=20),
            corner_radius=15,
            command=self.clear_search_history
        )
        self.btn_clear_history.pack(side="left", padx=(10, 0))

        self.btn_search = ctk.CTkButton(
            self.content_layer,
            text="🔍 开始搜索",
            height=50,
            command=self.perform_search_async,
            corner_radius=15,
            fg_color=COLORS["btn_primary"],
            hover_color=COLORS["btn_primary_hover"],
            text_color="white",
            font=ctk.CTkFont(size=18, weight="bold")
        )
        self.btn_search.pack(fill="x", pady=(0, 30), padx=40)

        self.results_header = ctk.CTkFrame(self.content_layer, height=40, corner_radius=10,
                                           fg_color=COLORS["header_bg"])
        self.results_header.pack(fill="x", padx=40, pady=(0, 0))
        ctk.CTkLabel(self.results_header, text="搜索结果", font=ctk.CTkFont(size=16, weight="bold"),
                     text_color=COLORS["header_text"]).place(relx=0.5, rely=0.5, anchor="center")

        self.results_scroll = ctk.CTkScrollableFrame(self.content_layer, label_text="", fg_color="transparent",
                                                     bg_color="transparent")
        self.results_scroll.pack(fill="both", expand=True, padx=40, pady=(10, 0))

        self.after(500, self.load_app_data)

    def _setup_background(self):
        texture_path = "paper_texture.jpg"
        if os.path.exists(texture_path):
            try:
                pil_img = Image.open(texture_path)
                self.bg_image_ctk = ctk.CTkImage(light_image=pil_img, dark_image=pil_img, size=(2560, 1440))
                self.bg_label = ctk.CTkLabel(self.main_container, text="", image=self.bg_image_ctk)
                self.bg_label.place(x=0, y=0, relwidth=1, relheight=1)
                self.content_layer.pack(fill="both", expand=True, pady=30)
                return
            except:
                pass
        self.content_layer.configure(fg_color=COLORS["bg_main"])
        self.content_layer.pack(fill="both", expand=True, pady=30)

    # 逻辑功能
    def load_app_data(self):
        if os.path.exists("history_search.json"):
            try:
                with open("history_search.json", "r", encoding='utf-8') as f:
                    self.search_history = json.load(f)
                    self.search_combo.configure(values=self.search_history)
            except:
                pass
        if os.path.exists("history_folders.json"):
            try:
                with open("history_folders.json", "r", encoding='utf-8') as f:
                    self.folder_history = json.load(f)
                    if self.folder_history:
                        self.folder_menu.configure(values=self.folder_history)
                        self.on_folder_change(self.folder_history[0])
                    else:
                        self.folder_menu.configure(values=["无历史记录"])
            except:
                pass

    def save_app_data(self):
        try:
            with open("history_search.json", "w", encoding='utf-8') as f:
                json.dump(self.search_history, f)
            with open("history_folders.json", "w", encoding='utf-8') as f:
                json.dump(self.folder_history, f)
        except:
            pass

    def get_index_path(self, folder_path):
        hash_name = hashlib.md5(folder_path.encode('utf-8')).hexdigest()
        return os.path.join(self.indexes_dir, f"index_{hash_name}.pkl")

    def browse_new_folder(self):
        path = filedialog.askdirectory()
        if path:
            if path not in self.folder_history:
                self.folder_history.insert(0, path)
                self.folder_menu.configure(values=self.folder_history)
                self.save_app_data()
            self.folder_var.set(path)
            self.on_folder_change(path)

    def on_folder_change(self, selected_folder):
        if selected_folder in ["无历史记录", "请选择或添加..."]: return
        self.current_folder = selected_folder
        self.folder_var.set(os.path.basename(selected_folder))
        self.status_label.configure(text=f"📂 选中库:\n{selected_folder}")
        index_file = self.get_index_path(selected_folder)
        if os.path.exists(index_file):
            success = self.engine.load_index_from_disk(index_file)
            if success:
                self.status_label.configure(text=f"☑ 已加载索引\n包含 {self.engine.total_docs} 篇文档")
                self.btn_search.configure(state="normal")
            else:
                self.status_label.configure(text="⚠️ 索引损坏，请重建")
        else:
            self.status_label.configure(text="⚠️ 此库无索引\n请点击下方按钮重建")
            self.engine = RankedSearchEngine()
        self.btn_index.configure(state="normal")

    def clear_search_history(self):
        if not self.search_history: return
        if messagebox.askyesno("确认", "确定要清空所有搜索记录吗？"):
            self.search_history = []
            self.search_combo.configure(values=[])
            self.search_combo.set("")
            self.save_app_data()

    def add_to_search_history(self, query):
        if not query: return
        if query in self.search_history: self.search_history.remove(query)
        self.search_history.insert(0, query)
        if len(self.search_history) > 15: self.search_history = self.search_history[:15]
        self.search_combo.configure(values=self.search_history)
        self.save_app_data()

    def start_indexing(self):
        self._set_ui_busy_state(True, "正在扫描并建立索引...\n请留意控制台输出")
        self.executor.submit(self.run_indexing_task)

    def run_indexing_task(self):
        count = 0
        try:
            print(f"--- 开始扫描文件夹: {self.current_folder} ---")
            # 重新初始化引擎，但传入当前的停用词表路径（如果有的话）
            self.engine = RankedSearchEngine()
            self.engine.indexed_folder = self.current_folder

            for root, dirs, files in os.walk(self.current_folder):
                for file in files:
                    full_path = os.path.join(root, file)
                    print(f"正在读取: {file}")
                    content = self.engine._extract_content(full_path)

                    if content.strip():
                        # 这里调用 add_document，现在它一定存在了
                        self.engine.add_document(count, content, full_path, file)
                        count += 1

            print(f"--- 扫描结束，共有效索引 {count} 个文件 ---")
            save_path = self.get_index_path(self.current_folder)
            self.engine.save_index_to_disk(save_path)
            self.after(0, lambda: self.finish_indexing(count))

        except Exception as e:
            err_msg = str(e)  # 【关键修复】立即捕获错误信息字符串
            print(f"× 索引过程错误: {err_msg}")
            # 使用 captured error string
            self.after(0, lambda: self._set_ui_busy_state(False, f"错误: {err_msg}"))

    def finish_indexing(self, count):
        self._set_ui_busy_state(False, f"√ 索引完成\n库中共有 {count} 篇文档")
        messagebox.showinfo("成功", f"索引建立完成！共 {count} 个有效文件。")

    def perform_search_async(self):
        query = self.search_combo.get().strip()
        if not query: return
        self.add_to_search_history(query)
        self.btn_search.configure(state="disabled", text="搜索中...")

        print(f"--- 开始搜索关键词: [{query}] ---")

        for widget in self.results_scroll.winfo_children(): widget.destroy()
        self.executor.submit(self.run_search_task, query)

    def run_search_task(self, query):
        try:
            results = self.engine.search(query)
            print(f"搜索完成，找到 {len(results)} 个结果")
            self.after(0, lambda: self.update_results_ui(results, query))
        except Exception as e:
            print(f"❌ 搜索过程出错: {e}")
            self.after(0, lambda: self._set_ui_busy_state(False, "搜索出错"))

    def update_results_ui(self, results, query):
        self.btn_search.configure(state="normal", text="🔍 开始搜索")
        if not results:
            ctk.CTkLabel(self.results_scroll, text="无搜索结果...", font=("Arial", 16),
                         text_color=COLORS["text_content"]).pack(pady=40)
            return
        for res in results:
            self.create_result_card(res, query)

    def create_result_card(self, res, query):
        def on_click(event):
            self.open_file(res['path'])

        card = ctk.CTkFrame(self.results_scroll, fg_color=COLORS["card_bg"], border_width=2,
                            border_color=COLORS["card_border"], corner_radius=15)
        card.pack(fill="x", padx=5, pady=8)
        card.bind("<Button-1>", on_click)
        card.bind("<Enter>", lambda e: card.configure(fg_color=COLORS["card_hover"]))
        card.bind("<Leave>", lambda e: card.configure(fg_color=COLORS["card_bg"]))

        title_frame = ctk.CTkFrame(card, fg_color="transparent")
        title_frame.pack(fill="x", padx=15, pady=(15, 5))
        title_frame.bind("<Button-1>", on_click)

        ctk.CTkLabel(title_frame, text=f"📄 {res['title']}", font=ctk.CTkFont(size=18, weight="bold"),
                     text_color=COLORS["text_main"]).pack(side="left")
        ctk.CTkLabel(title_frame, text=f"{res['score']}%", font=ctk.CTkFont(size=14, weight="bold"),
                     text_color="white", fg_color=COLORS["btn_primary"], corner_radius=8, width=50).pack(side="right",
                                                                                                         padx=5)
        ctk.CTkLabel(card, text=res['path'], font=ctk.CTkFont(size=11), text_color="gray", anchor="w").pack(fill="x",
                                                                                                            padx=15,
                                                                                                            pady=(
                                                                                                            0, 10))

        content = res['content'].replace('\n', ' ')
        idx = content.find(query)
        start = max(0, idx - 30)
        end = min(len(content), idx + 120)
        preview_text = "..." + content[start:end] + "..."

        preview_frame = ctk.CTkFrame(card, fg_color="transparent", corner_radius=8)
        preview_frame.pack(fill="x", padx=15, pady=(0, 15))
        parts = preview_text.split(query)
        for i, part in enumerate(parts):
            ctk.CTkLabel(preview_frame, text=part, text_color=COLORS["text_content"], font=ctk.CTkFont(size=13)).pack(
                side="left")
            if i < len(parts) - 1:
                ctk.CTkLabel(preview_frame, text=query, text_color=COLORS["text_highlight"],
                             font=ctk.CTkFont(size=13, weight="bold")).pack(side="left")
        ctk.CTkLabel(card, text="Click to open", font=ctk.CTkFont(size=10), text_color="#B0B0B0").pack(anchor="e",
                                                                                                       padx=15,
                                                                                                       pady=(0, 10))

    def open_file(self, path):
        try:
            os.startfile(path)
        except Exception as e:
            messagebox.showerror("Error", f"Cannot open file:\n{e}")

    def _set_ui_busy_state(self, is_busy, status_text):
        state = "disabled" if is_busy else "normal"
        self.btn_index.configure(state=state)
        self.btn_search.configure(state=state)
        self.btn_add_folder.configure(state=state)
        self.status_label.configure(text=status_text)
        if is_busy:
            self.progress.pack(padx=25, pady=10, fill="x")
            self.progress.start()
        else:
            self.progress.stop()
            self.progress.pack_forget()


if __name__ == "__main__":
    app = VibrantSearchApp()


    def on_closing():
        app.executor.shutdown(wait=False)
        app.destroy()


    app.protocol("WM_DELETE_WINDOW", on_closing)
    app.mainloop()
